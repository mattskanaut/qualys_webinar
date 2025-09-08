#!/usr/bin/env python3
"""
Create PowerPoint presentation for Episode 1: Demystifying the CNAPP
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import os

def create_title_slide(prs):
    """Create the title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Demystifying the CNAPP"
    subtitle.text = "Applying Vulnerability Management in the Cloud\n\nEpisode 1: Foundations"
    
    # Add notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Welcome everyone, let's cut through the confusion around cloud security"
    
    return slide

def create_agenda_slide(prs):
    """Create the agenda roadmap slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Our Journey Today"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "5 Big Questions We'll Answer:"
    tf.paragraphs[0].font.bold = True
    
    questions = [
        "What is 'best practice' in cloud?",
        "CNAPP - Another acronym or real solution?",
        "Why is cloud VM harder than on-prem?",
        "How do I get visibility into risk?",
        "Where to start Monday morning?"
    ]
    
    for i, q in enumerate(questions, 1):
        p = tf.add_paragraph()
        p.text = f"{i}. {q}"
        p.level = 0
        p.font.size = Pt(20)
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Five big questions we'll answer today - by the end you'll have practical actions"
    
    return slide

def create_question_slide(prs, question_num, question_text):
    """Create a question title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    title = slide.shapes.title
    title.text = f"Question {question_num}"
    
    # Add content
    content_shape = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8), Inches(3)
    )
    tf = content_shape.text_frame
    tf.text = question_text
    tf.word_wrap = True
    
    # Format text
    for paragraph in tf.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(32)
        paragraph.font.bold = True
    
    return slide

def create_breach_slide(prs, breach_name, details):
    """Create a breach story slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = f"Breach Story: {breach_name}"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    # Add breach details
    tf.text = details['overview']
    
    # Add timeline/impact
    p = tf.add_paragraph()
    p.text = f"\nDate: {details['date']}"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = f"Scale: {details['scale']}"
    
    # Add what went wrong
    p = tf.add_paragraph()
    p.text = "\nWhat Went Wrong:"
    p.font.bold = True
    
    for point in details['failures']:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.level = 1
    
    return slide

def create_capital_one_technical(prs):
    """Create the Capital One technical diagram slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Capital One: The Attack Path"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Attack Chain:"
    tf.paragraphs[0].font.bold = True
    
    steps = [
        "1. Internet → WAF (Misconfigured)",
        "2. SSRF Attack → Metadata Service",
        "3. Steal IAM Role Credentials",
        "4. Access S3 Buckets",
        "5. 100M Records Exposed"
    ]
    
    for step in steps:
        p = tf.add_paragraph()
        p.text = step
        p.font.size = Pt(24)
        if "Exposed" in step:
            p.font.bold = True
    
    # Add key insights
    p = tf.add_paragraph()
    p.text = "\nKey Failures:"
    p.font.bold = True
    
    failures = [
        "• No IMDSv2 enforcement",
        "• Over-permissioned IAM role",
        "• No network segmentation",
        "• Lack of anomaly detection"
    ]
    
    for failure in failures:
        p = tf.add_paragraph()
        p.text = failure
        p.font.size = Pt(20)
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Every piece was a small misconfiguration - together they created a highway to customer data"
    
    return slide

def create_shared_responsibility_slide(prs):
    """Create the shared responsibility model slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Shared Responsibility: The Reality"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    tf.text = "What Microsoft/AWS/GCP Secures:"
    tf.paragraphs[0].font.bold = True
    
    provider_items = [
        "• Physical infrastructure",
        "• Hypervisor",
        "• Network backbone"
    ]
    
    for item in provider_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "\nWhat YOU Secure:"
    p.font.bold = True
    
    customer_items = [
        "• Data",
        "• Applications", 
        "• Identity & Access",
        "• Operating Systems",
        "• Network & Firewall Config",
        "• Encryption",
        "• Everything else..."
    ]
    
    for item in customer_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
    
    # Add Gartner stat
    p = tf.add_paragraph()
    p.text = "\n⚠️ Gartner: 99% of cloud breaches will be customer's fault through 2025"
    p.font.bold = True
    p.font.size = Pt(22)
    
    return slide

def create_identity_perimeter_slide(prs):
    """Create the identity as perimeter slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Identity is the New Perimeter"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    tf.text = "Traditional Security:"
    tf.paragraphs[0].font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Castle & moat model"
    p = tf.add_paragraph()
    p.text = "• Trust inside the network"
    p = tf.add_paragraph()
    p.text = "• Firewall as main defense"
    
    p = tf.add_paragraph()
    p.text = "\nCloud Security:"
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• No perimeter"
    p = tf.add_paragraph()
    p.text = "• Zero trust architecture"
    p = tf.add_paragraph()
    p.text = "• Identity controls access"
    
    p = tf.add_paragraph()
    p.text = "\n💡 Key Insight:"
    p.font.bold = True
    p.font.size = Pt(24)
    
    p = tf.add_paragraph()
    p.text = "If someone gets an over-privileged managed identity or service principal, they don't need to 'hack in' - they just walk in."
    p.font.size = Pt(20)
    
    return slide

def create_poll_slide(prs, question, options):
    """Create a poll slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Quick Poll"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = question
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(28)
    
    p = tf.add_paragraph()
    p.text = ""  # Empty line
    
    for opt in options:
        p = tf.add_paragraph()
        p.text = opt
        p.font.size = Pt(24)
    
    return slide

def create_tool_sprawl_slide(prs):
    """Create the tool sprawl slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "The Tool Sprawl Problem"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Average Enterprise Security Stack:"
    tf.paragraphs[0].font.bold = True
    
    tools = [
        "• SIEM for logs",
        "• CSPM for misconfigurations",
        "• Vulnerability scanner for CVEs",
        "• CWPP for workloads",
        "• CIEM for identity risks",
        "• Container scanner",
        "• Secret scanner",
        "• IaC scanner",
        "• API security tool",
        "• Cloud native tools (Security Center, GuardDuty)",
        "• And 10+ more..."
    ]
    
    for tool in tools:
        p = tf.add_paragraph()
        p.text = tool
        p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "\n= 23 tools, 23 logins, 23 alert streams 😵"
    p.font.bold = True
    p.font.size = Pt(24)
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Each tool might be good, but they don't talk to each other"
    
    return slide

def create_cnapp_solution_slide(prs):
    """Create the CNAPP solution slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "CNAPP: The Unified Platform"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Cloud-Native Application Protection Platform"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(24)
    
    p = tf.add_paragraph()
    p.text = "\nWhat CNAPP Consolidates:"
    p.font.bold = True
    
    components = [
        "✓ CSPM - Cloud Security Posture Management",
        "✓ CWPP - Cloud Workload Protection",
        "✓ CIEM - Cloud Identity Entitlement Management",
        "✓ IaC Security - Infrastructure as Code scanning",
        "✓ Container & Kubernetes security",
        "✓ API discovery and protection"
    ]
    
    for comp in components:
        p = tf.add_paragraph()
        p.text = comp
        p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "\n🎯 Key Difference:"
    p.font.bold = True
    p.font.size = Pt(22)
    
    p = tf.add_paragraph()
    p.text = "Shared data model = Connected insights = Context-aware prioritization"
    p.font.size = Pt(20)
    
    return slide

def create_container_lifecycle_slide(prs):
    """Create the container lifecycle slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Why Cloud VM is Harder: Ephemeral Infrastructure"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Average Lifespans:"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(28)
    
    lifespans = [
        ("Traditional Server", "3-5 years"),
        ("Virtual Machine", "30-90 days"),
        ("Container", "12 minutes"),
        ("Serverless Function", "100 milliseconds"),
    ]
    
    for item, duration in lifespans:
        p = tf.add_paragraph()
        p.text = f"\n{item}:"
        p.font.bold = True
        p.font.size = Pt(24)
        
        p = tf.add_paragraph()
        p.text = f"    {duration}"
        p.font.size = Pt(22)
    
    p = tf.add_paragraph()
    p.text = "\n⚠️ Your 2 AM scan is outdated by 2:15 AM"
    p.font.bold = True
    p.font.size = Pt(24)
    
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "It's like trying to count cars on a highway - they keep moving"
    
    return slide

def create_visibility_matrix_slide(prs):
    """Create the visibility gap matrix slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Scanning Coverage Gaps"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "What Different Approaches Actually See:\n"
    tf.paragraphs[0].font.bold = True
    
    # Create a simple text-based table
    table_data = [
        "                        Agent    Network    API-Based",
        "Long-lived VMs:          ✓         ✓           ✓",
        "Containers:              ~         ✗           ✓",
        "Container Images:        ✗         ✗           ✓",
        "Serverless:              ✗         ✗           ✓",
        "PaaS Services:           ✗         ✗           ✓"
    ]
    
    for row in table_data:
        p = tf.add_paragraph()
        p.text = row
        p.font.name = "Courier New"
        p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "\n📊 If using agents alone: 60% blind spots"
    p.font.bold = True
    p.font.size = Pt(24)
    
    return slide

def create_risk_context_slide(prs):
    """Create the risk in context slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Context Changes Everything"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Same Vulnerability, Different Risk:"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(26)
    
    scenarios = [
        ("Dev Environment", "Internal only, no sensitive data", "LOW RISK ✅"),
        ("Production + WAF", "Protected, but customer-facing", "MEDIUM RISK ⚠️"),
        ("Internet-facing + Admin", "Exposed with privileged access", "CRITICAL RISK 🔴")
    ]
    
    for env, desc, risk in scenarios:
        p = tf.add_paragraph()
        p.text = f"\n{env}:"
        p.font.bold = True
        p.font.size = Pt(22)
        
        p = tf.add_paragraph()
        p.text = f"  {desc}"
        p.font.size = Pt(20)
        
        p = tf.add_paragraph()
        p.text = f"  → {risk}"
        p.font.bold = True
        p.font.size = Pt(22)
    
    p = tf.add_paragraph()
    p.text = "\n💡 CNAPP correlates all context for true risk scoring"
    p.font.size = Pt(20)
    
    return slide

def create_30_day_roadmap_slide(prs):
    """Create the 30-60-90 day roadmap slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Your 30-60-90 Day Roadmap"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    roadmap = [
        ("Week 1: Discovery", [
            "Enable Azure Activity Log",
            "Turn on Security Center",
            "Run asset discovery",
            "Document findings"
        ]),
        ("Week 2-4: Quick Wins", [
            "Close public storage accounts",
            "Enable MFA everywhere",
            "Rotate old keys",
            "Tag resources"
        ]),
        ("Month 2-3: Operationalize", [
            "Deploy automated scanning",
            "Implement risk prioritization",
            "Train the team",
            "Build remediation workflows"
        ])
    ]
    
    for phase, tasks in roadmap:
        p = tf.add_paragraph()
        p.text = phase
        p.font.bold = True
        p.font.size = Pt(24)
        
        for task in tasks:
            p = tf.add_paragraph()
            p.text = f"  • {task}"
            p.font.size = Pt(18)
    
    return slide

def create_success_story_slide(prs):
    """Create the success story slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Real Success Story"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "50-person Fintech, All-in on Azure"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(24)
    
    p = tf.add_paragraph()
    p.text = "\nBefore:"
    p.font.bold = True
    p.font.size = Pt(22)
    
    before = [
        "• 2000+ vulnerabilities",
        "• 400+ misconfigurations",
        "• No AKS visibility",
        "• 3-person team drowning"
    ]
    
    for item in before:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "\nAfter 90 Days:"
    p.font.bold = True
    p.font.size = Pt(22)
    
    after = [
        "• 50 critical vulns (prioritized)",
        "• Automated remediation",
        "• Full container scanning",
        "• Same 3 people, now strategic"
    ]
    
    for item in after:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "\n🎯 Key: They didn't fix everything - they fixed what mattered"
    p.font.bold = True
    p.font.size = Pt(20)
    
    return slide

def create_key_takeaways_slide(prs):
    """Create the key takeaways slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Takeaways"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    takeaways = [
        ("🔄", "Cloud security isn't harder, it's different"),
        ("👁️", "Visibility must be continuous, not periodic"),
        ("🎯", "Context matters more than severity"),
        ("🚀", "Start small, win fast, then scale")
    ]
    
    for emoji, text in takeaways:
        p = tf.add_paragraph()
        p.text = f"{emoji} {text}"
        p.font.size = Pt(26)
        p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "\n\nRemember: Perfect security is impossible, but good enough security is achievable"
    p.font.size = Pt(20)
    
    return slide

def create_next_steps_slide(prs):
    """Create the next steps slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Next Steps"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    tf.text = "Continue Your Journey:"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(28)
    
    steps = [
        "\n📅 Episode 2: Deep Dive into Cloud-Native VM",
        "   Next week - Technical implementation details",
        "\n📚 Resources:",
        "   • Cloud Security Best Practices Guide",
        "   • CNAPP Evaluation Checklist",
        "   • Azure Security Benchmark",
        "\n🎯 Take Action:",
        "   • Book a demo",
        "   • Join our community Slack",
        "   • Start your 30-day roadmap"
    ]
    
    for step in steps:
        p = tf.add_paragraph()
        p.text = step
        p.font.size = Pt(20)
        if step.startswith("\n"):
            p.font.bold = True
    
    return slide

def create_thank_you_slide(prs):
    """Create the thank you/Q&A slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Thank You!"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    
    tf.text = "Q&A Time"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(36)
    
    p = tf.add_paragraph()
    p.text = "\n\nDrop your questions in the chat!"
    p.font.size = Pt(24)
    
    p = tf.add_paragraph()
    p.text = "\n\nContact:"
    p.font.bold = True
    p.font.size = Pt(20)
    
    p = tf.add_paragraph()
    p.text = "Email: [your-email]"
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "LinkedIn: [your-linkedin]"
    p.font.size = Pt(18)
    
    p = tf.add_paragraph()
    p.text = "Qualys.com/cloud-security"
    p.font.size = Pt(18)
    
    return slide

def main():
    """Create the complete presentation"""
    # Create presentation
    prs = Presentation()
    
    # Set slide size to 16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Create slides
    print("Creating Episode 1 presentation...")
    
    # Opening
    create_title_slide(prs)
    create_agenda_slide(prs)
    
    # Question 1: Best Practices
    create_question_slide(prs, 1, "What does 'security best practice'\nactually mean in a cloud-native world?")
    
    # Capital One breach
    capital_one = {
        'date': 'March-July 2019',
        'scale': '106 million customers affected',
        'overview': 'A former AWS employee exploited a misconfigured WAF to access customer data via SSRF attack.',
        'failures': [
            'Misconfigured Web Application Firewall',
            'Over-permissioned IAM role',
            'No IMDSv2 enforcement',
            'Lack of network segmentation',
            'No anomaly detection for S3 access'
        ]
    }
    create_breach_slide(prs, "Capital One (2019)", capital_one)
    create_capital_one_technical(prs)
    create_shared_responsibility_slide(prs)
    create_identity_perimeter_slide(prs)
    
    # Poll 1
    create_poll_slide(
        prs,
        "What surprised you most about the Capital One breach?",
        [
            "A. The simplicity of the attack",
            "B. The IAM role issues",
            "C. That it happened to a tech-savvy bank",
            "D. The metadata service vector"
        ]
    )
    
    # Question 2: CNAPP
    create_question_slide(prs, 2, "CNAPP sounds like another acronym.\nWhat problem is it actually solving?")
    create_tool_sprawl_slide(prs)
    
    # Uber breach
    uber = {
        'date': 'October 2016',
        'scale': '57 million users, 600,000 driver licenses',
        'overview': 'Attackers found credentials on GitHub, accessed AWS S3 buckets with customer data.',
        'failures': [
            'Credentials in code repository',
            'No MFA on cloud accounts',
            'No secrets scanning',
            'Fragmented security tools',
            'Tool blind spots between GitHub and AWS'
        ]
    }
    create_breach_slide(prs, "Uber (2016)", uber)
    create_cnapp_solution_slide(prs)
    
    # Question 3: Cloud VM Challenges
    create_question_slide(prs, 3, "Why is vulnerability management\nso much harder in the cloud?")
    create_container_lifecycle_slide(prs)
    
    # Tesla breach
    tesla = {
        'date': 'February 2018',
        'scale': 'Unknown data exposure, cryptomining impact',
        'overview': 'Kubernetes dashboard exposed without password, led to cryptomining and data access.',
        'failures': [
            'Kubernetes dashboard without authentication',
            'Default settings not hardened',
            'AWS credentials in pods',
            'No network policies',
            'No runtime monitoring'
        ]
    }
    create_breach_slide(prs, "Tesla Kubernetes (2018)", tesla)
    create_visibility_matrix_slide(prs)
    create_risk_context_slide(prs)
    
    # Question 4: Visibility
    create_question_slide(prs, 4, "How do I get visibility\ninto my actual cloud risk?")
    
    # SolarWinds breach
    solarwinds = {
        'date': 'March-December 2020',
        'scale': '18,000 organizations affected',
        'overview': 'Nation-state actors compromised build environment, distributed backdoored updates.',
        'failures': [
            'Build environment compromised',
            'Weak password (solarwinds123)',
            'Code signing abuse',
            'Supply chain trust exploited',
            'Months of undetected presence'
        ]
    }
    create_breach_slide(prs, "SolarWinds (2020)", solarwinds)
    
    # Question 5: Where to Start
    create_question_slide(prs, 5, "Where should a team start\nwith cloud VM?")
    create_30_day_roadmap_slide(prs)
    create_success_story_slide(prs)
    
    # Closing
    create_key_takeaways_slide(prs)
    create_next_steps_slide(prs)
    create_thank_you_slide(prs)
    
    # Save presentation
    filename = 'Episode1_CloudSecurity.pptx'
    prs.save(filename)
    print(f"✅ Presentation saved as {filename}")
    print(f"   Total slides: {len(prs.slides)}")
    
    return filename

if __name__ == "__main__":
    main()